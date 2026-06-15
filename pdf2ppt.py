#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF 转 PPT 脚本
从 PDF 提取内容，生成美化版 PPT
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pdfplumber import PDF

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_DEEP = RGBColor(0x04, 0x0b, 0x18)
BG_DARK = RGBColor(0x06, 0x11, 0x1f)
BG_CARD = RGBColor(0x0c, 0x1f, 0x35)

CYAN = RGBColor(0x00, 0xe5, 0xff)
CYAN_DIM = RGBColor(0x00, 0xa8, 0xcc)
GOLD = RGBColor(0xff, 0xb3, 0x00)
BLUE = RGBColor(0x1a, 0x90, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xcc, 0xd6, 0xe0)
MUTED = RGBColor(0x6b, 0x7d, 0x8c)

FONT = "PingFang SC"
FONT_MONO = "Menlo"

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


def extract_pdf_content(pdf_path):
    """提取 PDF 内容"""
    content = {
        "pages": [],
        "total_pages": 0
    }

    with open(pdf_path, "rb") as f:
        with PDF(f) as pdf:
            content["total_pages"] = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages):
                page_data = {
                    "page_num": page_num + 1,
                    "text": "",
                    "tables": [],
                    "images": []
                }

                text = page.extract_text()
                if text:
                    page_data["text"] = text.strip()

                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_data = []
                        for row in table:
                            if row:
                                table_data.append([cell.strip() if cell else "" for cell in row])
                        if table_data:
                            page_data["tables"].append(table_data)

                content["pages"].append(page_data)

    return content


def determine_slide_structure(content):
    """根据 PDF 内容确定 PPT 结构"""
    pages = content["pages"]
    total = content["total_pages"]

    structure = []
    for i in range(total):
        if i == 0:
            structure.append("cover")
        elif i == total - 1:
            structure.append("end")
        else:
            structure.append("content")
    return structure


def create_slide(prs, slide_type, page_data, page_num, total, structure_idx):
    """创建单张幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if slide_type == "cover":
        make_cover_slide(slide, page_data, page_num, total)
    elif slide_type == "end":
        make_end_slide(slide, page_data, page_num, total)
    elif slide_type == "section":
        make_section_slide(slide, page_data, page_num, total)
    else:
        make_content_slide(slide, page_data, page_num, total)

    return slide


def make_cover_slide(slide, page_data, page_num, total):
    """封面页"""
    set_slide_bg(slide, BG_DEEP)

    add_top_bar(slide)
    add_corner_deco(slide)

    text = page_data.get("text", "")
    lines = text.split("\n") if text else []

    main_title = ""
    subtitle = ""
    other_lines = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if not main_title and len(line) < 30:
            main_title = line
        elif not subtitle and len(line) < 40:
            subtitle = line
        else:
            other_lines.append(line)

    if not main_title:
        main_title = "智慧低空应急运输教学平台"
    if not subtitle:
        subtitle = "Smart Low-Altitude Emergency Transportation"

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = main_title
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(24)
        r2.font.color.rgb = CYAN
        r2.font.name = FONT
        p2.alignment = PP_ALIGN.LEFT

    if other_lines:
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        for i, line in enumerate(other_lines[:3]):
            p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(16)
            r.font.color.rgb = LIGHT
            r.font.name = FONT

    tags = ["教学平台", "AI 智能体", "岗课赛证", "低空应急运输"]
    tag_y = Inches(6.5)
    for i, tag in enumerate(tags):
        tag_box = slide.shapes.add_textbox(Inches(1) + i * Inches(2.8), tag_y, Inches(2.5), Inches(0.4))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tag
        r.font.size = Pt(14)
        r.font.bold = True
        colors = [CYAN, GOLD, RGBColor(0x00, 0xd4, 0xaa), RGBColor(0xa8, 0x55, 0xf7)]
        r.font.color.rgb = colors[i % len(colors)]
        r.font.name = FONT_MONO

    add_footer(slide, page_num, total)


def make_content_slide(slide, page_data, page_num, total):
    """内容页"""
    set_slide_bg(slide, BG_DARK)
    add_top_bar(slide)

    text = page_data.get("text", "")
    lines = text.split("\n") if text else []

    title = ""
    body_lines = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if not title and len(line) < 25:
            title = line
        else:
            body_lines.append(line)

    if not title:
        title = f"第 {page_num} 页"

    add_title_with_bar(slide, title, "", Inches(0.8))

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), SLIDE_W - Inches(1.6), SLIDE_H - Inches(3))
    ctf = content_box.text_frame
    ctf.word_wrap = True

    if body_lines:
        chunk_size = 5
        for i in range(0, len(body_lines), chunk_size):
            chunk = body_lines[i:i+chunk_size]
            for j, line in enumerate(chunk):
                p = ctf.paragraphs[0] if i == 0 and j == 0 else ctf.add_paragraph()
                p.space_before = Pt(6)
                r = p.add_run()
                r.text = line
                r.font.size = Pt(16)
                r.font.color.rgb = LIGHT
                r.font.name = FONT

    tables = page_data.get("tables", [])
    if tables:
        for t_idx, table in enumerate(tables[:2]):
            rows = len(table)
            cols = len(table[0]) if table else 1
            graphic_frame = slide.shapes.add_table(
                rows, cols,
                Inches(0.8), Inches(4.5 + t_idx * 1.5),
                SLIDE_W - Inches(1.6), Inches(1.2)
            )
            table_obj = graphic_frame.table
            for row_idx, row in enumerate(table):
                for col_idx, cell in enumerate(row):
                    cell_obj = table_obj.cell(row_idx, col_idx)
                    cell_obj.text = str(cell)[:30]
                    for paragraph in cell_obj.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(11)
                            run.font.color.rgb = WHITE if row_idx == 0 else LIGHT
                            run.font.bold = (row_idx == 0)

    add_footer(slide, page_num, total)


def make_section_slide(slide, page_data, page_num, total):
    """章节分隔页"""
    set_slide_bg(slide, BG_DEEP)
    add_top_bar(slide)
    add_corner_deco(slide)

    text = page_data.get("text", "")
    lines = text.split("\n") if text else []

    title = f"第 {page_num} 节"
    for line in lines:
        if line.strip() and len(line.strip()) < 20:
            title = line.strip()
            break

    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W/2 - Inches(1.5), Inches(2.5), Inches(3), Inches(0.04))
    deco.line.fill.background()
    deco.fill.solid()
    deco.fill.fore_color.rgb = CYAN

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.0), SLIDE_W - Inches(2), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), SLIDE_W - Inches(2), Inches(1))
    stf = sub_box.text_frame
    p2 = stf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"第 {page_num} 页 / 共 {total} 页"
    r2.font.size = Pt(18)
    r2.font.color.rgb = CYAN
    r2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER

    add_footer(slide, page_num, total)


def make_end_slide(slide, page_data, page_num, total):
    """结尾页"""
    set_slide_bg(slide, BG_DEEP)
    add_top_bar(slide)
    add_corner_deco(slide)

    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W/2 - Inches(2), Inches(2.0), Inches(4), Inches(0.03))
    deco.line.fill.background()
    deco.fill.solid()
    deco.fill.fore_color.rgb = CYAN

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), SLIDE_W - Inches(2), Inches(2))
    tf = thanks_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "感谢聆听"
    r.font.size = Pt(72)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), SLIDE_W - Inches(2), Inches(2))
    stf = sub_box.text_frame
    stf.word_wrap = True

    info = [
        "🚁 智慧低空应急运输教学平台",
        "⚙️ AI × 无人机 × 应急物流",
    ]
    for i, line in enumerate(info):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.space_before = Pt(12)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18)
        r.font.color.rgb = LIGHT
        r.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, page_num, total)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.05))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE


def add_corner_deco(slide):
    for x, y in [(0, 0), (SLIDE_W, 0), (0, SLIDE_H), (SLIDE_W, SLIDE_H)]:
        corner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.4), Inches(0.02))
        corner.line.fill.background()
        corner.fill.solid()
        corner.fill.fore_color.rgb = CYAN


def add_title_with_bar(slide, title, subtitle, top=Inches(0.8)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.8), Inches(0.04))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN

    tb = slide.shapes.add_textbox(Inches(0.8), top + Inches(0.15), SLIDE_W - Inches(1.6), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = CYAN
        r2.font.name = FONT


def add_footer(slide, n, total):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), SLIDE_H - Inches(0.6), SLIDE_W - Inches(1.2), Inches(0.01))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1a, 0x30, 0x50)

    left_tb = slide.shapes.add_textbox(Inches(0.6), SLIDE_H - Inches(0.55), Inches(6), Inches(0.4))
    ltf = left_tb.text_frame
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = "智慧低空应急运输教学平台"
    lr.font.size = Pt(11)
    lr.font.color.rgb = MUTED
    lr.font.name = FONT

    num_tb = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.55), Inches(0.8), Inches(0.4))
    ntf = num_tb.text_frame
    np = ntf.paragraphs[0]
    nr = np.add_run()
    nr.text = f"{n} / {total}"
    nr.font.size = Pt(12)
    nr.font.color.rgb = CYAN
    nr.font.name = FONT_MONO
    nr.font.bold = True
    np.alignment = PP_ALIGN.RIGHT


def main():
    pdf_path = os.path.join(OUT_DIR, "智慧低空应急运输教学平台.pdf")

    if not os.path.exists(pdf_path):
        print(f"❌ PDF file not found: {pdf_path}")
        sys.exit(1)

    print(f"📄 Reading PDF: {pdf_path}")
    content = extract_pdf_content(pdf_path)
    print(f"📊 Found {content['total_pages']} pages")

    print("📐 Determining slide structure...")
    structure = determine_slide_structure(content)
    total_slides = len(structure)
    print(f"🎯 Will create {total_slides} slides: {structure}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, slide_type in enumerate(structure):
        page_idx = min(idx, len(content["pages"]) - 1)
        page_data = content["pages"][page_idx]

        print(f"  📝 Slide {idx+1}/{total_slides}: {slide_type} (from PDF page {page_idx+1})")
        create_slide(prs, slide_type, page_data, idx + 1, total_slides, idx)

    out_path = os.path.join(OUT_DIR, "智慧低空应急运输教学平台_PPT版.pptx")
    prs.save(out_path)
    print(f"✅ Saved: {out_path}")
    print(f"📏 Size: {os.path.getsize(out_path)} bytes")


if __name__ == "__main__":
    main()
