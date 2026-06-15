from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_DEEP = RGBColor(0x04, 0x0b, 0x18)
BG_DARK = RGBColor(0x06, 0x11, 0x1f)
BG_CARD = RGBColor(0x0c, 0x1f, 0x35)

CYAN = RGBColor(0x00, 0xe5, 0xff)
CYAN_DIM = RGBColor(0x00, 0xa8, 0xcc)
GOLD = RGBColor(0xff, 0xb3, 0x00)
GOLD_DIM = RGBColor(0xcc, 0x8f, 0x00)
BLUE = RGBColor(0x1a, 0x90, 0xff)
PURPLE = RGBColor(0xa8, 0x55, 0xf7)
GREEN = RGBColor(0x00, 0xd4, 0xaa)
RED_ACCENT = RGBColor(0xff, 0x47, 0x57)

WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xcc, 0xd6, 0xe0)
MUTED = RGBColor(0x6b, 0x7d, 0x8c)

FONT = "PingFang SC"
FONT_MONO = "Menlo"

OUT = os.path.join(os.path.dirname(__file__), "低空应急智能运输_路演PPT.pptx")


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_decorative_corners(slide, color=CYAN):
    for (x, y) in [(0, 0), (SLIDE_W, 0), (0, SLIDE_H), (SLIDE_W, SLIDE_H)]:
        corner = slide.shapes.add_shape(MSO_SHAPE.LINE, x, y, Inches(0.4), Inches(0.02))
        corner.line.color.rgb = color
        corner.line.width = Pt(1.5)
        corner.fill.background()


def add_top_accent(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.05))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.05), SLIDE_W, Inches(0.02))
    glow.line.fill.background()
    glow.fill.solid()
    glow.fill.fore_color.rgb = CYAN_DIM


def add_footer(slide, n, total):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), SLIDE_H - Inches(0.6), SLIDE_W - Inches(1.2), Inches(0.01))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1a, 0x30, 0x50)

    tb = slide.shapes.add_textbox(Inches(0.6), SLIDE_H - Inches(0.55), Inches(8), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "智慧低空应急运输教学平台"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED
    r.font.name = FONT

    num = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.55), Inches(0.8), Inches(0.4))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"{n} / {total}"
    r2.font.size = Pt(12)
    r2.font.color.rgb = CYAN
    r2.font.name = FONT_MONO
    r2.font.bold = True
    p2.alignment = PP_ALIGN.RIGHT

    dots = slide.shapes.add_textbox(SLIDE_W / 2 - Inches(1), SLIDE_H - Inches(0.55), Inches(2), Inches(0.4))
    tf3 = dots.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    parts = []
    for i in range(total):
        if i < n - 1:
            parts.append("●")
        elif i == n - 1:
            parts.append("●")
        else:
            parts.append("○")
    r3 = p3.add_run()
    r3.text = "  ".join(parts)
    r3.font.size = Pt(8)
    r3.font.color.rgb = CYAN if n else MUTED
    r3.font.name = FONT_MONO
    p3.alignment = PP_ALIGN.CENTER


def add_title_block(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.7) + Inches(0.5), Inches(0.8), Inches(0.04))
    accent_line.line.fill.background()
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = CYAN

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = LIGHT
        r2.font.name = FONT


def add_glow_card(slide, left, top, width, height, title, body=None, icon="", accent=CYAN):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    card.shadow.inherit = False
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD

    glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Inches(0.05), top - Inches(0.05), width + Inches(0.1), height + Inches(0.1))
    glow.line.color.rgb = accent
    glow.line.width = Pt(0.5)
    glow.line.transparency = 0.7
    glow.fill.background()

    accent_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    accent_strip.line.fill.background()
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = accent

    content_tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = content_tb.text_frame
    tf.word_wrap = True

    if icon:
        icon_tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.12), Inches(0.5), Inches(0.4))
        icon_tf = icon_tb.text_frame
        ip = icon_tf.paragraphs[0]
        ir = ip.add_run()
        ir.text = icon
        ir.font.size = Pt(24)

    title_tb = slide.shapes.add_textbox(left + Inches(0.85) if icon else left + Inches(0.3), top + Inches(0.12), width - (Inches(0.9) if icon else Inches(0.4)), Inches(0.45))
    title_tf = title_tb.text_frame
    tp = title_tf.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(17)
    tr.font.bold = True
    tr.font.color.rgb = WHITE
    tr.font.name = FONT

    if body:
        body_tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.65), width - Inches(0.4), height - Inches(0.8))
        body_tf = body_tb.text_frame
        body_tf.word_wrap = True
        for i, line in enumerate(body):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.space_before = Pt(2)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(13)
            r.font.color.rgb = LIGHT
            r.font.name = FONT


def slide1_cover(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DEEP)

    add_top_accent(s)
    add_decorative_corners(s, CYAN_DIM)
    add_footer(s, 1, total)

    deco1 = s.shapes.add_shape(MSO_SHAPE.LINE, Inches(1), Inches(1.2), Inches(3), Inches(0.01))
    deco1.line.color.rgb = CYAN
    deco1.line.width = Pt(1)
    deco1.line.transparency = 0.5

    deco2 = s.shapes.add_shape(MSO_SHAPE.LINE, Inches(1), Inches(1.5), Inches(2), Inches(0.01))
    deco2.line.color.rgb = CYAN
    deco2.line.width = Pt(0.5)
    deco2.line.transparency = 0.7

    main_title = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2.5))
    tf = main_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "智慧低空应急运输"
    r.font.size = Pt(64)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "教学平台"
    r2.font.size = Pt(56)
    r2.font.bold = True
    r2.font.color.rgb = CYAN
    r2.font.name = FONT

    tag_tb = s.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(0.6))
    tag_tf = tag_tb.text_frame
    tp = tag_tf.paragraphs[0]
    tr = tp.add_run()
    tr.text = "🚁  AI × 无人机 × 应急物流"
    tr.font.size = Pt(22)
    tr.font.color.rgb = GOLD
    tr.font.name = FONT

    desc_tb = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    desc_tf = desc_tb.text_frame
    desc_tf.word_wrap = True
    lines = [
        "用工程化的路径规划，让低空运输不再是纸上谈兵",
        "一个平台 + 三个AI智能体 · 面向高职/本科教学实训",
    ]
    for i, line in enumerate(lines):
        p = desc_tf.paragraphs[0] if i == 0 else desc_tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(16)
        r.font.color.rgb = LIGHT
        r.font.name = FONT

    side_deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(2), Inches(2.5), Inches(0.04), Inches(3))
    side_deco.line.fill.background()
    side_deco.fill.solid()
    side_deco.fill.fore_color.rgb = CYAN_DIM

    tags = [
        ("教学平台", CYAN),
        ("AI 智能体", GOLD),
        ("岗课赛证", GREEN),
        ("案例:渠洋村", PURPLE),
    ]
    tag_start_x = Inches(1)
    tag_y = Inches(6.5)
    for i, (text, color) in enumerate(tags):
        tb = s.shapes.add_textbox(tag_start_x + i * Inches(2.5), tag_y, Inches(2.3), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = FONT_MONO


def slide2_problem(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DARK)
    add_top_accent(s)
    add_decorative_corners(s, CYAN_DIM)
    add_footer(s, 2, total)
    add_title_block(s, "痛点 · 为什么要做", "三个真实问题，驱动我们的项目")

    items = [
        (
            "🚁",
            "应急物流卡脖子",
            "灾区路毁、桥梁断\n地面运输彻底瘫痪\n低空成唯一生命线",
            RED_ACCENT,
        ),
        (
            "📈",
            "人才严重缺口",
            "2025低空经济破2万亿\n无人机物流岗爆发\n学校缺实训平台",
            GOLD,
        ),
        (
            "🎯",
            "教学脱离实际",
            "只教理论没有约束\n学生不知载重几趟\nAI当实训助理改作业",
            CYAN,
        ),
    ]

    card_w = Inches(3.8)
    card_h = Inches(3.8)
    gap = Inches(0.5)
    start_x = (SLIDE_W - card_w * 3 - gap * 2) / 2
    top_y = Inches(2.2)

    for i, (icon, title, body, accent) in enumerate(items):
        add_glow_card(
            s,
            start_x + i * (card_w + gap),
            top_y,
            card_w,
            card_h,
            title,
            body.split("\n"),
            icon,
            accent,
        )

    bottom = s.shapes.add_textbox(Inches(1), Inches(6.4), SLIDE_W - Inches(2), Inches(0.6))
    tf = bottom.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "→ 解决方案：一套能跑、能学、能评的教学实训平台"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = CYAN
    r.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


def slide3_solution(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DARK)
    add_top_accent(s)
    add_decorative_corners(s, CYAN_DIM)
    add_footer(s, 3, total)
    add_title_block(s, "项目全貌", "1 个教学平台 + 3 个 AI 智能体")

    platform_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5))
    platform_card.line.color.rgb = CYAN
    platform_card.line.width = Pt(2)
    platform_card.fill.solid()
    platform_card.fill.fore_color.rgb = BG_CARD

    icon_tb = s.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(0.6), Inches(0.5))
    icon_tf = icon_tb.text_frame
    ip = icon_tf.paragraphs[0]
    ir = ip.add_run()
    ir.text = "🏛️"
    ir.font.size = Pt(28)

    title_tb = s.shapes.add_textbox(Inches(1.6), Inches(2.3), Inches(4.5), Inches(0.5))
    title_tf = title_tb.text_frame
    tp = title_tf.paragraphs[0]
    tr = tp.add_run()
    tr.text = "教学平台"
    tr.font.size = Pt(22)
    tr.font.bold = True
    tr.font.color.rgb = CYAN
    tr.font.name = FONT

    features = [
        "首页驾驶舱 · 数据大屏",
        "课程中心 · 实训发布",
        "四主体教学智评",
        "学习资源库",
        "AI 助教浮窗",
    ]
    for i, feat in enumerate(features):
        y_pos = Inches(2.9) + i * Inches(0.5)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y_pos + Inches(0.08), Inches(0.08), Inches(0.08))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN

        feat_tb = s.shapes.add_textbox(Inches(1.3), y_pos, Inches(4.8), Inches(0.4))
        feat_tf = feat_tb.text_frame
        fp = feat_tf.paragraphs[0]
        fr = fp.add_run()
        fr.text = feat
        fr.font.size = Pt(15)
        fr.font.color.rgb = LIGHT
        fr.font.name = FONT

    arrow_shape = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), Inches(3.8), Inches(0.5), Inches(0.3))
    arrow_shape.line.color.rgb = CYAN
    arrow_shape.line.width = Pt(1)
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = CYAN_DIM

    agents = [
        (
            "🧭",
            "路径规划智能体",
            "蚁群 CVRP 算法\n动态能耗模型\n渠洋村 8 村案例",
            CYAN,
        ),
        (
            "📦",
            "装箱评价智能体",
            "空间利用率\n重量平衡分析\n安全评分",
            GOLD,
        ),
        (
            "📚",
            "课程图谱智能体",
            "知识/能力/问题/思政\n学习路径推荐\n教学关联分析",
            GREEN,
        ),
    ]

    agent_w = Inches(2.2)
    agent_h = Inches(2.0)
    agent_gap = Inches(0.3)
    agent_start_x = Inches(7.1)
    agent_top_y = Inches(2.1)

    for i, (icon, title, body, accent) in enumerate(agents):
        add_glow_card(
            s,
            agent_start_x + i * (agent_w + agent_gap),
            agent_top_y,
            agent_w,
            agent_h,
            title,
            body.split("\n"),
            icon,
            accent,
        )

    connector = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(5.2), Inches(6.6), Inches(0.01))
    connector.line.color.rgb = CYAN_DIM
    connector.line.width = Pt(0.5)
    connector.line.transparency = 0.5

    api_tb = s.shapes.add_textbox(Inches(6.4), Inches(5.3), Inches(6.6), Inches(0.5))
    api_tf = api_tb.text_frame
    ap = api_tf.paragraphs[0]
    ar = ap.add_run()
    ar.text = "API 接入 · 零耦合 · 独立开发部署"
    ar.font.size = Pt(12)
    ar.font.color.rgb = MUTED
    ar.font.name = FONT_MONO
    ap.alignment = PP_ALIGN.CENTER

    bottom = s.shapes.add_textbox(Inches(1), Inches(6.2), SLIDE_W - Inches(2), Inches(0.6))
    tf = bottom.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "→ 学生登录平台 → 选择实训任务 → 调用智能体 → AI 评分 → 形成能力画像"
    r.font.size = Pt(16)
    r.font.color.rgb = LIGHT
    r.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


def slide4_core_tech(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DARK)
    add_top_accent(s)
    add_decorative_corners(s, CYAN_DIM)
    add_footer(s, 4, total)
    add_title_block(s, "核心技术亮点", "工程化可运行的 CVRP 路径规划引擎")

    algo_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), SLIDE_W - Inches(1.6), Inches(2.2))
    algo_card.line.color.rgb = CYAN
    algo_card.line.width = Pt(1.5)
    algo_card.fill.solid()
    algo_card.fill.fore_color.rgb = BG_CARD

    algo_title = s.shapes.add_textbox(Inches(1.1), Inches(2.15), Inches(3), Inches(0.45))
    algo_tf = algo_title.text_frame
    ap = algo_tf.paragraphs[0]
    ar = ap.add_run()
    ar.text = "🐜 蚁群算法 CVRP"
    ar.font.size = Pt(18)
    ar.font.bold = True
    ar.font.color.rgb = WHITE
    ar.font.name = FONT

    formula = s.shapes.add_textbox(Inches(1.1), Inches(2.65), SLIDE_W - Inches(2.2), Inches(0.6))
    formula_tf = formula.text_frame
    fp = formula_tf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "目标函数"
    fr.font.size = Pt(11)
    fr.font.color.rgb = MUTED
    fr.font.name = FONT

    formula2 = s.shapes.add_textbox(Inches(1.1), Inches(2.95), SLIDE_W - Inches(2.2), Inches(0.5))
    f2_tf = formula2.text_frame
    f2p = f2_tf.paragraphs[0]

    segments = [
        ("0.4", "距离", CYAN),
        (" + 0.3", "能耗", GOLD),
        (" + 0.2", "优先级", RED_ACCENT),
        (" + 0.1", "负载均衡", GREEN),
    ]
    for i, (weight, label, color) in enumerate(segments):
        wr = f2p.add_run()
        wr.text = weight
        wr.font.size = Pt(18)
        wr.font.bold = True
        wr.font.color.rgb = color
        wr.font.name = FONT_MONO

        lr = f2p.add_run()
        lr.text = f"·{label}"
        lr.font.size = Pt(14)
        lr.font.color.rgb = LIGHT
        lr.font.name = FONT
        if i < len(segments) - 1:
            pr = f2p.add_run()
            pr.text = "  "
            pr.font.size = Pt(14)

    tech_cards = [
        (
            "⚡",
            "动态能耗模型",
            [
                "能耗 = 距离 × (1 + 载重/最大载重)",
                "满载去程 = 空载返航的 2 倍",
                "载重越高，能耗越大",
            ],
            GOLD,
        ),
        (
            "📏",
            "载重-航程插值",
            [
                "0kg → 26km（空载）",
                "80kg → 6km（满载）",
                "分段线性插值计算实时航程",
            ],
            CYAN,
        ),
        (
            "🤖",
            "AI 任务规划 Agent",
            [
                "大模型自动选机型",
                "重量/航程/抗风/冷链 多维评估",
                "输出方案 + 风险预警",
            ],
            PURPLE,
        ),
    ]

    tech_w = Inches(3.8)
    tech_h = Inches(2.5)
    tech_gap = Inches(0.4)
    tech_start_x = (SLIDE_W - tech_w * 3 - tech_gap * 2) / 2
    tech_top_y = Inches(4.5)

    for i, (icon, title, body, accent) in enumerate(tech_cards):
        add_glow_card(
            s,
            tech_start_x + i * (tech_w + tech_gap),
            tech_top_y,
            tech_w,
            tech_h,
            title,
            body,
            icon,
            accent,
        )


def slide5_value(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DARK)
    add_top_accent(s)
    add_decorative_corners(s, CYAN_DIM)
    add_footer(s, 5, total)
    add_title_block(s, "价值主张", "三个维度，让项目真正落地")

    values = [
        (
            "🎓",
            "对学生",
            [
                "亲手规划救灾航线",
                "理解载重-航程非线性关系",
                "产出就业作品集",
                "准备 1+X 证书",
            ],
            CYAN,
        ),
        (
            "👨‍🏫",
            "对老师",
            [
                "真实案例 + 真实数据",
                "AI 助教 7×24 答疑",
                "四主体教学闭环",
                "创新创业项目底座",
            ],
            GOLD,
        ),
        (
            "🏢",
            "对专业",
            [
                "低空经济/应急/物流 三方向",
                "可扩展：装箱/集群/禁飞区",
                "可申请教科研课题",
                "企业专家可接入",
            ],
            GREEN,
        ),
    ]

    val_w = Inches(4.0)
    val_h = Inches(4.0)
    val_gap = Inches(0.4)
    val_start_x = (SLIDE_W - val_w * 3 - val_gap * 2) / 2
    val_top_y = Inches(2.1)

    for i, (icon, title, body, accent) in enumerate(values):
        add_glow_card(
            s,
            val_start_x + i * (val_w + val_gap),
            val_top_y,
            val_w,
            val_h,
            title,
            body,
            icon,
            accent,
        )

    vision = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.3), SLIDE_W - Inches(3), Inches(0.7))
    vision.line.color.rgb = CYAN
    vision.line.width = Pt(1)
    vision.fill.solid()
    vision.fill.fore_color.rgb = RGBColor(0x0a, 0x20, 0x40)

    v_tb = s.shapes.add_textbox(Inches(1.7), Inches(6.35), SLIDE_W - Inches(3.4), Inches(0.6))
    v_tf = v_tb.text_frame
    vp = v_tf.paragraphs[0]
    vr = vp.add_run()
    vr.text = "💡  让每一个想做无人机物流的学生，有一台「随时可练、随时可评」的教学飞行塔"
    vr.font.size = Pt(16)
    vr.font.bold = True
    vr.font.color.rgb = CYAN
    vr.font.name = FONT
    vp.alignment = PP_ALIGN.CENTER


def slide6_end(prs, total):
    s = blank_slide(prs)
    set_bg(s, BG_DEEP)
    add_top_accent(s)
    add_decorative_corners(s, CYAN)
    add_footer(s, 6, total)

    deco_line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W / 2 - Inches(1.5), Inches(1.8), Inches(3), Inches(0.03))
    deco_line.line.fill.background()
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = CYAN

    deco_line2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W / 2 - Inches(1), Inches(1.9), Inches(2), Inches(0.01))
    deco_line2.line.fill.background()
    deco_line2.fill.solid()
    deco_line2.fill.fore_color.rgb = CYAN_DIM

    main = s.shapes.add_textbox(Inches(1), Inches(2.2), SLIDE_W - Inches(2), Inches(1.5))
    tf = main.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "让低空物流"
    r.font.size = Pt(64)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "飞进每一个课堂"
    r2.font.size = Pt(56)
    r2.font.bold = True
    r2.font.color.rgb = CYAN
    r2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER

    deco_line3 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W / 2 - Inches(2), Inches(4.0), Inches(4), Inches(0.02))
    deco_line3.line.fill.background()
    deco_line3.fill.solid()
    deco_line3.fill.fore_color.rgb = CYAN_DIM

    details = s.shapes.add_textbox(Inches(1), Inches(4.2), SLIDE_W - Inches(2), Inches(2))
    dtf = details.text_frame
    dtf.word_wrap = True

    info = [
        "🚁  案例：广西渠洋村 8 村应急物资配送",
        "⚙️  技术栈：Vue 3 · Django · Python 蚁群算法 · LLM",
        "🎯  5 分钟路演 · 创业训练项目",
    ]
    for i, line in enumerate(info):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        p.space_before = Pt(12)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18)
        r.font.color.rgb = LIGHT
        r.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    thanks = s.shapes.add_textbox(Inches(1), Inches(6.3), SLIDE_W - Inches(2), Inches(0.8))
    ttf = thanks.text_frame
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = "感 谢 聆 听"
    tr.font.size = Pt(36)
    tr.font.bold = True
    tr.font.color.rgb = GOLD
    tr.font.name = FONT
    tp.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 6
    slide1_cover(prs, total)
    slide2_problem(prs, total)
    slide3_solution(prs, total)
    slide4_core_tech(prs, total)
    slide5_value(prs, total)
    slide6_end(prs, total)

    prs.save(OUT)
    print("✅ Saved:", OUT)


if __name__ == "__main__":
    main()
